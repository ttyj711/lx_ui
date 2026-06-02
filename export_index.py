#!/usr/bin/env python3
"""
Export index.html as an editable PPTX file.
Applies footer-rail + cursor-flow discipline from pptx-html-fidelity-audit skill.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Canvas (16:9)
CANVAS_W = Inches(13.333)
CANVAS_H = Inches(7.5)

# Margins
MARGIN_X = Inches(0.6)
MARGIN_TOP = Inches(0.5)
CONTENT_LEFT = MARGIN_X
CONTENT_RIGHT = CANVAS_W - MARGIN_X
CONTENT_W = CONTENT_RIGHT - CONTENT_LEFT

# Vertical rails
CONTENT_MAX_Y = Inches(6.70)
FOOTER_TOP = Inches(6.85)
FOOTER_H = Inches(0.22)

# Theme colors from index.html
COLOR_BG = RGBColor(0xF5, 0xF5, 0xF7)       # --bg
COLOR_SURFACE = RGBColor(0xFF, 0xFF, 0xFF)  # --surface
COLOR_FG = RGBColor(0x1D, 0x1D, 0x1F)       # --fg
COLOR_MUTED = RGBColor(0xA0, 0xA4, 0xAE)    # --muted
COLOR_PRIMARY = RGBColor(0x2D, 0x8B, 0x7A)  # --primary
COLOR_PRIMARY_LIGHT = RGBColor(0xE6, 0xF3, 0xF0)  # --primary-light
COLOR_PRICE = RGBColor(0xE0, 0x5A, 0x38)    # --price
COLOR_PRICE_LIGHT = RGBColor(0xFD, 0xF0, 0xEC)  # --price-light

# Typography
FONT_DISPLAY = "PingFang SC"
FONT_BODY = "PingFang SC"

class Cursor:
    """Advances down the slide; refuses to cross the footer rail."""
    def __init__(self, y_start=MARGIN_TOP, cap=CONTENT_MAX_Y):
        self.y = y_start
        self.cap = cap
        self.history = []

    def take(self, h, gap=Inches(0.12), label=""):
        top = self.y
        self.y = top + h + gap
        self.history.append((top, h, label))
        if self.y > self.cap:
            raise OverflowError(
                f"Cursor exceeded rail at '{label}': y={self.y} cap={self.cap}"
            )
        return top

def hero_layout(blocks):
    """Center a stack of blocks vertically on the canvas."""
    total_h = sum(h + g for h, g in blocks)
    y_start = (CANVAS_H - total_h) / 2
    return Cursor(y_start=y_start, cap=CANVAS_H - FOOTER_H - Inches(0.2))

def set_slide_background(slide, color):
    """Set slide background to a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, *, left, top, width, height, text, font, size_pt,
             color=None, bold=False, italic=False, align="left", name=None):
    """Add a text box with specified formatting."""
    shape = slide.shapes.add_textbox(left, top, width, height)
    if name:
        shape.name = name
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]

    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic and not _is_cjk(text)
    if color:
        run.font.color.rgb = color

    # Set CJK font explicitly
    if _is_cjk(text):
        run.font.name = font
        try:
            run._r.get_or_add_rPr()
            ea = run._r.rPr.find(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}ea'
            )
            if ea is None:
                from lxml import etree
                ea = etree.SubElement(
                    run._r.rPr,
                    '{http://schemas.openxmlformats.org/drawingml/2006/main}ea'
                )
            ea.set('typeface', font)
        except Exception:
            pass

    return shape

def _is_cjk(s):
    """Check if string contains CJK characters."""
    return any('一' <= c <= '鿿' or '぀' <= c <= 'ヿ' for c in s)

def add_rectangle(slide, *, left, top, width, height, fill_color=None, radius=Inches(0), name=None):
    """Add a rectangle shape with optional fill."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if name:
        shape.name = name
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()

    # Set corner radius
    if radius > Inches(0):
        w_in = width.inches if hasattr(width, 'inches') else width / 914400
        h_in = height.inches if hasattr(height, 'inches') else height / 914400
        r_in = radius.inches if hasattr(radius, 'inches') else radius / 914400
        shape.adjustments[0] = min(0.5, r_in / min(w_in, h_in))

    return shape

def add_footer(slide, text, theme="light"):
    """Add footer at FOOTER_TOP."""
    color = COLOR_MUTED
    add_text(slide,
        left=CONTENT_LEFT, top=FOOTER_TOP,
        width=CONTENT_W, height=FOOTER_H,
        text=text, font=FONT_BODY, size_pt=12,
        color=color, align="center", name="footer")

def create_pptx():
    """Create the PPTX presentation."""
    prs = Presentation()
    prs.slide_width = CANVAS_W
    prs.slide_height = CANVAS_H

    blank_layout = prs.slide_layouts[6]

    # Slide 1: Cover/Hero
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, COLOR_PRIMARY)

    blocks = [
        (Inches(0.35), Inches(0.30)),   # badge
        (Inches(0.80), Inches(0.20)),   # title
        (Inches(0.90), Inches(0.00)),   # description
    ]
    c = hero_layout(blocks)

    # Badge
    badge_top = c.take(*blocks[0])
    add_rectangle(slide1,
        left=Inches(5.2), top=badge_top,
        width=Inches(2.9), height=Inches(0.35),
        fill_color=RGBColor(0x19, 0x62, 0x53), radius=Inches(0.12))

    add_text(slide1,
        left=Inches(5.2), top=badge_top + Inches(0.05),
        width=Inches(2.9), height=Inches(0.25),
        text="📱 移动端 APP 原型",
        font=FONT_BODY, size_pt=13,
        color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, align="center")

    # Title
    add_text(slide1,
        left=CONTENT_LEFT, top=c.take(*blocks[1]),
        width=CONTENT_W, height=Inches(0.80),
        text="旅途账本",
        font=FONT_DISPLAY, size_pt=56,
        color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, align="center")

    # Description
    add_text(slide1,
        left=Inches(2.5), top=c.take(*blocks[2]),
        width=Inches(8.3), height=Inches(0.90),
        text="多人旅行费用记录与分摊 APP — 快速记账、智能分摊、旅行总结,让每一段旅程的账目清清楚楚。",
        font=FONT_BODY, size_pt=18,
        color=RGBColor(0xE8, 0xF3, 0xF0), align="center")

    add_footer(slide1, "旅途账本 · Open Design 原型", theme="dark")

    # Slide 2: Features
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, COLOR_BG)

    c = Cursor()

    add_text(slide2,
        left=CONTENT_LEFT, top=c.take(Inches(0.60), label="title"),
        width=CONTENT_W, height=Inches(0.60),
        text="核心功能",
        font=FONT_DISPLAY, size_pt=36,
        color=COLOR_FG, bold=True, align="center")

    c.take(Inches(0.30))

    features = [
        ("一键记账", "快速录入消费", COLOR_PRIMARY, COLOR_PRIMARY_LIGHT),
        ("多人分摊", "智能费用分配", COLOR_PRICE, COLOR_PRICE_LIGHT),
        ("智能统计", "消费结构分析", RGBColor(0xE6, 0x51, 0x00), RGBColor(0xFF, 0xF3, 0xE0)),
        ("旅行总结", "账单导出分享", RGBColor(0x15, 0x65, 0xC0), RGBColor(0xE3, 0xF2, 0xFD)),
    ]

    card_w = (CONTENT_W - Inches(0.9)) / 4
    card_h = Inches(2.4)
    row_top = c.take(card_h, label="features")

    for i, (title, desc, icon_color, bg_color) in enumerate(features):
        left = CONTENT_LEFT + i * (card_w + Inches(0.3))

        add_rectangle(slide2,
            left=left, top=row_top,
            width=card_w, height=card_h,
            fill_color=COLOR_SURFACE, radius=Inches(0.18))

        icon_size = Inches(0.72)
        add_rectangle(slide2,
            left=left + (card_w - icon_size) / 2,
            top=row_top + Inches(0.4),
            width=icon_size, height=icon_size,
            fill_color=bg_color, radius=Inches(0.18))

        add_text(slide2,
            left=left, top=row_top + Inches(1.3),
            width=card_w, height=Inches(0.40),
            text=title,
            font=FONT_DISPLAY, size_pt=18,
            color=COLOR_FG, bold=True, align="center")

        add_text(slide2,
            left=left, top=row_top + Inches(1.75),
            width=card_w, height=Inches(0.35),
            text=desc,
            font=FONT_BODY, size_pt=12,
            color=COLOR_MUTED, align="center")

    add_footer(slide2, "旅途账本 · Open Design 原型")

    # Slide 3: Gallery Overview
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, COLOR_BG)

    c = Cursor()

    add_text(slide3,
        left=CONTENT_LEFT, top=c.take(Inches(0.50), label="title"),
        width=CONTENT_W, height=Inches(0.50),
        text="界面总览",
        font=FONT_DISPLAY, size_pt=32,
        color=COLOR_FG, bold=True, align="center")

    add_text(slide3,
        left=CONTENT_LEFT, top=c.take(Inches(0.30), label="subtitle"),
        width=CONTENT_W, height=Inches(0.30),
        text="5个核心界面,覆盖完整记账流程",
        font=FONT_BODY, size_pt=14,
        color=COLOR_MUTED, align="center")

    c.take(Inches(0.40))

    screens = [
        ("首页 · 旅程概况", "预算进度、今日消费、通知提醒"),
        ("记一笔 · 快速录入", "分类选择、参与成员、小票上传"),
        ("统计 · 消费分析", "消费结构、每日趋势、人均排行"),
        ("总结 · 旅行报告", "消费总结、亮点分析、导出分享"),
        ("设置 · 配置中心", "旅程默认、成员分摊、个性化配置"),
    ]

    screen_w = (CONTENT_W - Inches(0.8)) / 5
    screen_h = Inches(3.5)
    row_top = c.take(screen_h, label="screens")

    for i, (title, desc) in enumerate(screens):
        left = CONTENT_LEFT + i * (screen_w + Inches(0.2))

        add_rectangle(slide3,
            left=left, top=row_top,
            width=screen_w, height=screen_h,
            fill_color=COLOR_SURFACE, radius=Inches(0.12))

        add_text(slide3,
            left=left, top=row_top + Inches(0.3),
            width=screen_w, height=Inches(0.35),
            text=title,
            font=FONT_DISPLAY, size_pt=14,
            color=COLOR_FG, bold=True, align="center")

        add_text(slide3,
            left=left, top=row_top + Inches(0.65),
            width=screen_w, height=Inches(0.60),
            text=desc,
            font=FONT_BODY, size_pt=11,
            color=COLOR_MUTED, align="center")

    add_footer(slide3, "旅途账本 · Open Design 原型")

    # Slide 4: Closing
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, COLOR_PRIMARY)

    blocks = [
        (Inches(0.80), Inches(0.20)),
        (Inches(0.90), Inches(0.00)),
    ]
    c = hero_layout(blocks)

    add_text(slide4,
        left=CONTENT_LEFT, top=c.take(*blocks[0]),
        width=CONTENT_W, height=Inches(0.80),
        text="让旅途账目清清楚楚",
        font=FONT_DISPLAY, size_pt=42,
        color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, align="center")

    add_text(slide4,
        left=CONTENT_LEFT, top=c.take(*blocks[1]),
        width=CONTENT_W, height=Inches(0.90),
        text="快速记账 · 智能分摊 · 轻松总结",
        font=FONT_BODY, size_pt=20,
        color=RGBColor(0xE8, 0xF3, 0xF0), align="center")

    add_footer(slide4, "旅途账本 · Open Design 原型", theme="dark")

    output_path = "index.pptx"
    prs.save(output_path)
    return output_path

if __name__ == "__main__":
    path = create_pptx()
    print(f"Created: {path}")
